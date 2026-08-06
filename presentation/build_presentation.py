"""
NovaCRM Enterprise Presentation Builder
Creates a 23-slide 16:9 PowerPoint using real application screenshots.
"""

from __future__ import annotations

import math
import os
from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "assets" / "screenshots"
ASSETS = ROOT / "presentation" / "assets"
OUT_DIR = ROOT / "presentation"
PPTX_PATH = OUT_DIR / "NovaCRM_Presentation.pptx"
PDF_PATH = OUT_DIR / "NovaCRM_Presentation.pdf"

# Brand
PRIMARY = RGBColor(0x25, 0x63, 0xEB)
SECONDARY = RGBColor(0x1E, 0x40, 0xAF)
ACCENT = RGBColor(0x0E, 0xA5, 0xE9)
DARK = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
CARD = RGBColor(0xF1, 0xF5, 0xF9)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)

W, H = Inches(13.333), Inches(7.5)


def ensure_dirs() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    SHOTS.mkdir(parents=True, exist_ok=True)


def hex_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return tuple(int(h[i : i + 2], 16) for i in (0, 2, 4))


def gradient_image(size: tuple[int, int], c1: str, c2: str, diagonal: bool = True) -> Image.Image:
    w, h = size
    img = Image.new("RGB", (w, h))
    r1, g1, b1 = hex_rgb(c1)
    r2, g2, b2 = hex_rgb(c2)
    px = img.load()
    for y in range(h):
        for x in range(w):
            t = ((x / max(w - 1, 1)) + (y / max(h - 1, 1))) / 2 if diagonal else y / max(h - 1, 1)
            px[x, y] = (
                int(r1 + (r2 - r1) * t),
                int(g1 + (g2 - g1) * t),
                int(b1 + (b2 - b1) * t),
            )
    return img


def make_cover_bg() -> Path:
    path = ASSETS / "cover-bg.png"
    img = gradient_image((1920, 1080), "#0B1B3A", "#2563EB")
    draw = ImageDraw.Draw(img, "RGBA")
    # Soft orbs
    for cx, cy, r, color in [
        (1500, 200, 320, (14, 165, 233, 70)),
        (300, 850, 280, (30, 64, 175, 90)),
        (1100, 700, 220, (37, 99, 235, 60)),
    ]:
        overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
        ImageDraw.Draw(overlay).ellipse((cx - r, cy - r, cx + r, cy + r), fill=color)
        img = Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")
    # Grid lines
    draw = ImageDraw.Draw(img, "RGBA")
    for x in range(0, 1920, 80):
        draw.line([(x, 0), (x, 1080)], fill=(255, 255, 255, 18), width=1)
    for y in range(0, 1080, 80):
        draw.line([(0, y), (1920, y)], fill=(255, 255, 255, 18), width=1)
    img.save(path, "PNG")
    return path


def make_abstract_panel(name: str, c1: str, c2: str) -> Path:
    path = ASSETS / f"{name}.png"
    img = gradient_image((1200, 800), c1, c2)
    draw = ImageDraw.Draw(img, "RGBA")
    for i in range(6):
        y = 80 + i * 110
        draw.rounded_rectangle((80, y, 520, y + 70), radius=16, fill=(255, 255, 255, 28))
    for i in range(4):
        y = 120 + i * 140
        draw.rounded_rectangle((600, y, 1120, y + 90), radius=18, fill=(255, 255, 255, 22))
    img.save(path, "PNG")
    return path


def round_screenshot(src: Path, dest: Path, radius: int = 28, shadow: bool = True) -> Path:
    if not src.exists():
        # placeholder card
        img = Image.new("RGB", (1440, 900), (241, 245, 249))
        d = ImageDraw.Draw(img)
        d.text((40, 40), f"Missing: {src.name}", fill=(100, 116, 139))
        src_img = img
    else:
        src_img = Image.open(src).convert("RGBA")

    # Crop to content if huge
    src_img = src_img.resize((1440, 900), Image.Resampling.LANCZOS)
    mask = Image.new("L", src_img.size, 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, *src_img.size), radius=radius, fill=255)
    rounded = Image.new("RGBA", src_img.size, (0, 0, 0, 0))
    rounded.paste(src_img, mask=mask)

    if shadow:
        canvas = Image.new("RGBA", (src_img.width + 40, src_img.height + 40), (0, 0, 0, 0))
        sh = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
        ImageDraw.Draw(sh).rounded_rectangle(
            (20, 24, 20 + src_img.width, 24 + src_img.height),
            radius=radius,
            fill=(15, 23, 42, 70),
        )
        sh = sh.filter(ImageFilter.GaussianBlur(14))
        canvas = Image.alpha_composite(canvas, sh)
        canvas.paste(rounded, (20, 16), rounded)
        canvas.save(dest, "PNG")
    else:
        rounded.save(dest, "PNG")
    return dest


def make_icon(name: str, color: str = "#2563EB", bg: str = "#EFF6FF") -> Path:
    path = ASSETS / f"icon-{name}.png"
    size = 256
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((8, 8, size - 8, size - 8), radius=56, fill=hex_rgb(bg) + (255,))
    c = hex_rgb(color)
    # Simple geometric glyphs
    cx, cy = size // 2, size // 2
    if name in {"disconnect", "systems"}:
        draw.ellipse((70, 70, 130, 130), outline=c, width=10)
        draw.ellipse((126, 126, 186, 186), outline=c, width=10)
    elif name in {"hr", "manual"}:
        draw.ellipse((98, 70, 158, 130), outline=c, width=10)
        draw.arc((70, 140, 186, 210), 20, 160, fill=c, width=10)
    elif name in {"excel", "sheets"}:
        draw.rectangle((78, 70, 178, 190), outline=c, width=10)
        for y in (100, 130, 160):
            draw.line((90, y, 166, y), fill=c, width=6)
    elif name in {"collab", "people"}:
        draw.ellipse((70, 90, 120, 140), outline=c, width=8)
        draw.ellipse((136, 90, 186, 140), outline=c, width=8)
        draw.arc((50, 150, 140, 210), 20, 160, fill=c, width=8)
        draw.arc((116, 150, 206, 210), 20, 160, fill=c, width=8)
    elif name in {"visibility", "eye"}:
        draw.ellipse((60, 100, 196, 160), outline=c, width=10)
        draw.ellipse((110, 112, 146, 148), fill=c)
    elif name in {"delay", "clock"}:
        draw.ellipse((70, 70, 186, 186), outline=c, width=10)
        draw.line((128, 128, 128, 90), fill=c, width=10)
        draw.line((128, 128, 160, 145), fill=c, width=10)
    elif name in {"cost", "money"}:
        draw.ellipse((78, 78, 178, 178), outline=c, width=10)
        draw.line((128, 95, 128, 161), fill=c, width=10)
    elif name in {"prod", "zap"}:
        draw.polygon([(140, 70), (100, 130), (130, 130), (116, 186), (170, 112), (138, 112)], fill=c)
    elif name in {"central", "hub"}:
        draw.ellipse((108, 108, 148, 148), fill=c)
        for ang in range(0, 360, 60):
            rad = math.radians(ang)
            x = cx + int(70 * math.cos(rad))
            y = cy + int(70 * math.sin(rad))
            draw.ellipse((x - 14, y - 14, x + 14, y + 14), outline=c, width=8)
            draw.line((cx, cy, x, y), fill=c, width=6)
    elif name in {"auto", "gear"}:
        draw.ellipse((100, 100, 156, 156), outline=c, width=10)
        for ang in range(0, 360, 45):
            rad = math.radians(ang)
            x1 = cx + int(40 * math.cos(rad))
            y1 = cy + int(40 * math.sin(rad))
            x2 = cx + int(70 * math.cos(rad))
            y2 = cy + int(70 * math.sin(rad))
            draw.line((x1, y1, x2, y2), fill=c, width=10)
    elif name in {"report", "chart"}:
        draw.rectangle((70, 70, 186, 186), outline=c, width=8)
        draw.rectangle((90, 140, 110, 170), fill=c)
        draw.rectangle((120, 110, 140, 170), fill=c)
        draw.rectangle((150, 90, 170, 170), fill=c)
    elif name in {"scale", "growth"}:
        draw.line((70, 180, 186, 80), fill=c, width=10)
        draw.polygon([(186, 80), (150, 88), (178, 116)], fill=c)
    elif name in {"security", "shield"}:
        draw.polygon([(128, 60), (190, 90), (190, 140), (128, 196), (66, 140), (66, 90)], outline=c, width=10)
        draw.line((105, 128, 122, 148), fill=c, width=10)
        draw.line((122, 148, 160, 108), fill=c, width=10)
    elif name in {"cloud"}:
        draw.ellipse((70, 110, 140, 170), fill=c)
        draw.ellipse((110, 90, 180, 160), fill=c)
        draw.ellipse((140, 120, 200, 170), fill=c)
        draw.rectangle((90, 135, 180, 170), fill=c)
    elif name in {"module"}:
        draw.rounded_rectangle((70, 70, 120, 120), radius=12, fill=c)
        draw.rounded_rectangle((136, 70, 186, 120), radius=12, fill=c)
        draw.rounded_rectangle((70, 136, 120, 186), radius=12, fill=c)
        draw.rounded_rectangle((136, 136, 186, 186), radius=12, fill=c)
    elif name in {"check"}:
        draw.ellipse((60, 60, 196, 196), outline=c, width=10)
        draw.line((90, 130, 118, 160), fill=c, width=12)
        draw.line((118, 160, 170, 100), fill=c, width=12)
    else:
        draw.ellipse((88, 88, 168, 168), outline=c, width=12)
    img.save(path, "PNG")
    return path


def make_qr_placeholder() -> Path:
    path = ASSETS / "qr.png"
    img = Image.new("RGB", (320, 320), "white")
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 319, 319), outline=(15, 23, 42), width=8)
    # Fake QR pattern
    for y in range(20, 300, 18):
        for x in range(20, 300, 18):
            if (x * y + x + y) % 47 < 18:
                draw.rectangle((x, y, x + 12, y + 12), fill=(15, 23, 42))
    for box in [(24, 24, 90, 90), (230, 24, 296, 90), (24, 230, 90, 296)]:
        draw.rectangle(box, outline=(15, 23, 42), width=6)
        draw.rectangle((box[0] + 16, box[1] + 16, box[2] - 16, box[3] - 16), fill=(15, 23, 42))
    img.save(path, "PNG")
    return path


def make_device_mockups() -> Path:
    path = ASSETS / "devices.png"
    canvas = Image.new("RGBA", (1600, 900), (248, 250, 252, 255))
    # laptop frame
    lap = Image.open(SHOTS / "01-dashboard.png").convert("RGBA").resize((900, 560), Image.Resampling.LANCZOS)
    frame = Image.new("RGBA", (960, 620), (0, 0, 0, 0))
    ImageDraw.Draw(frame).rounded_rectangle((0, 0, 959, 579), radius=18, fill=(15, 23, 42, 255))
    frame.paste(lap, (30, 20), lap)
    ImageDraw.Draw(frame).rounded_rectangle((180, 590, 780, 619), radius=8, fill=(30, 41, 59, 255))
    canvas.paste(frame, (40, 140), frame)

    # tablet
    tab = Image.open(SHOTS / "02-hrms-dashboard.png").convert("RGBA").resize((420, 560), Image.Resampling.LANCZOS)
    tframe = Image.new("RGBA", (460, 600), (0, 0, 0, 0))
    ImageDraw.Draw(tframe).rounded_rectangle((0, 0, 459, 599), radius=28, fill=(15, 23, 42, 255))
    tframe.paste(tab, (20, 20), tab)
    canvas.paste(tframe, (980, 140), tframe)

    # phone
    ph = Image.open(SHOTS / "12-tasks-board.png").convert("RGBA").crop((0, 0, 480, 900)).resize((220, 420), Image.Resampling.LANCZOS)
    pframe = Image.new("RGBA", (250, 460), (0, 0, 0, 0))
    ImageDraw.Draw(pframe).rounded_rectangle((0, 0, 249, 459), radius=36, fill=(15, 23, 42, 255))
    pframe.paste(ph, (15, 20), ph)
    canvas.paste(pframe, (1320, 220), pframe)

    canvas.save(path, "PNG")
    return path


def set_run(run, size=18, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_text(shape, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return tf


def add_para(tf, text, size=16, bold=False, color=SLATE, align=PP_ALIGN.LEFT, space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_bg(slide, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def title_block(slide, title: str, subtitle: str | None = None):
    add_accent_bar(slide)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.6))
    add_text(box, title, size=32, bold=True, color=DARK)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(12), Inches(0.4))
        add_text(sub, subtitle, size=16, color=MUTED)


def card(slide, left, top, width, height, fill=LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.adjustments[0] = 0.1
    return shape


def prepare_rounded_shots() -> dict[str, Path]:
    mapping = {}
    for src in SHOTS.glob("*.png"):
        dest = ASSETS / f"rounded-{src.name}"
        mapping[src.stem] = round_screenshot(src, dest)
    return mapping


def build() -> Path:
    ensure_dirs()
    cover_bg = make_cover_bg()
    panel = make_abstract_panel("ecosystem-panel", "#1E40AF", "#0EA5E9")
    qr = make_qr_placeholder()
    devices = make_device_mockups()
    shots = prepare_rounded_shots()

    icons = {
        k: make_icon(k)
        for k in [
            "disconnect",
            "manual",
            "excel",
            "collab",
            "visibility",
            "delay",
            "central",
            "prod",
            "report",
            "security",
            "module",
            "cloud",
            "cost",
            "auto",
            "scale",
            "check",
            "hub",
            "chart",
            "zap",
            "gear",
            "shield",
            "people",
            "eye",
            "clock",
            "money",
            "growth",
        ]
    }

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ---------- Slide 1 Cover ----------
    s = blank_slide(prs)
    s.shapes.add_picture(str(cover_bg), 0, 0, W, H)
    badge = card(s, Inches(0.8), Inches(1.6), Inches(2.2), Inches(0.45), RGBColor(0x1E, 0x3A, 0x8A))
    add_text(badge, "  ENTERPRISE PLATFORM", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(8), Inches(1.2))
    add_text(t, "NovaCRM", size=60, bold=True, color=WHITE)
    t2 = s.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8), Inches(0.5))
    add_text(t2, "Business Management Platform", size=24, color=RGBColor(0xBF, 0xDB, 0xFE))
    t3 = s.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(9), Inches(1))
    tf = add_text(t3, "One Platform.", size=28, bold=True, color=WHITE)
    add_para(tf, "Complete Business Control.", size=28, bold=True, color=WHITE, space_before=2)
    foot = s.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(10), Inches(0.4))
    add_text(foot, "Confidential  ·  Enterprise Sales Deck  ·  2026", size=12, color=RGBColor(0x93, 0xC5, 0xFD))

    # ---------- Slide 2 About ----------
    s = blank_slide(prs)
    add_bg(s)
    title_block(s, "About NovaCRM", "A unified operating system for modern enterprises")
    para = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(1.1))
    add_text(
        para,
        "NovaCRM brings HRMS, Projects, CRM, Finance readiness, Inventory, Helpdesk, and Automation "
        "into one secure workspace — replacing fragmented tools with a single source of truth for people, "
        "customers, and delivery.",
        size=16,
        color=SLATE,
    )
    modules = [
        ("HRMS", PRIMARY),
        ("Projects", SECONDARY),
        ("CRM", ACCENT),
        ("Finance", GREEN),
        ("Inventory", AMBER),
        ("Helpdesk", ROSE),
        ("Automation", RGBColor(0x8B, 0x5C, 0xF6)),
    ]
    x = 0.6
    for i, (label, color) in enumerate(modules):
        c = card(s, Inches(x), Inches(3.0), Inches(1.55), Inches(2.4), LIGHT)
        c.fill.fore_color.rgb = color
        tb = s.shapes.add_textbox(Inches(x), Inches(3.9), Inches(1.55), Inches(0.8))
        add_text(tb, label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(modules) - 1:
            arrow = s.shapes.add_textbox(Inches(x + 1.45), Inches(3.9), Inches(0.3), Inches(0.5))
            add_text(arrow, "→", size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        x += 1.8

    # ---------- Slide 3 Problems ----------
    s = blank_slide(prs)
    add_bg(s)
    title_block(s, "Business Problems", "What fragmented operations cost every growing company")
    problems = [
        ("disconnect", "Disconnected Software"),
        ("manual", "Manual HR"),
        ("excel", "Excel Sheets"),
        ("collab", "Poor Collaboration"),
        ("visibility", "No Visibility"),
        ("delay", "Delayed Decisions"),
    ]
    positions = [
        (0.6, 1.6),
        (4.7, 1.6),
        (8.8, 1.6),
        (0.6, 4.3),
        (4.7, 4.3),
        (8.8, 4.3),
    ]
    for (icon, label), (lx, ty) in zip(problems, positions):
        card(s, Inches(lx), Inches(ty), Inches(3.7), Inches(2.3), LIGHT)
        s.shapes.add_picture(str(icons[icon]), Inches(lx + 1.3), Inches(ty + 0.35), Inches(1.0), Inches(1.0))
        tb = s.shapes.add_textbox(Inches(lx + 0.2), Inches(ty + 1.5), Inches(3.3), Inches(0.6))
        add_text(tb, label, size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # ---------- Slide 4 Why ----------
    s = blank_slide(prs)
    add_bg(s)
    title_block(s, "Why NovaCRM", "Outcomes that matter to leadership teams")
    reasons = [
        ("central", "Centralized Operations", "One workspace for people, projects, and revenue"),
        ("prod", "Higher Productivity", "Less switching, faster execution, clearer ownership"),
        ("report", "Real-Time Analytics", "Live KPIs for attendance, delivery, and sales"),
        ("security", "Enterprise Security", "RBAC, audit logs, encryption-ready controls"),
        ("module", "Modular Architecture", "Adopt modules as your organization scales"),
        ("cloud", "Cloud Ready", "Accessible anywhere with modern web architecture"),
    ]
    for i, (icon, title, desc) in enumerate(reasons):
        col, row = i % 3, i // 3
        lx, ty = 0.6 + col * 4.2, 1.6 + row * 2.7
        card(s, Inches(lx), Inches(ty), Inches(3.95), Inches(2.4), LIGHT)
        s.shapes.add_picture(str(icons[icon]), Inches(lx + 0.3), Inches(ty + 0.35), Inches(0.7), Inches(0.7))
        tb = s.shapes.add_textbox(Inches(lx + 0.3), Inches(ty + 1.2), Inches(3.35), Inches(0.4))
        add_text(tb, "✓  " + title, size=16, bold=True, color=DARK)
        db = s.shapes.add_textbox(Inches(lx + 0.3), Inches(ty + 1.65), Inches(3.35), Inches(0.55))
        add_text(db, desc, size=13, color=MUTED)

    # ---------- Slide 5 Ecosystem ----------
    s = blank_slide(prs)
    add_bg(s)
    title_block(s, "Platform Overview", "Everything connected. One enterprise graph.")
    hub = card(s, Inches(5.1), Inches(3.1), Inches(3.1), Inches(1.2), PRIMARY)
    add_text(hub, "NovaCRM", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    nodes = [
        (1.0, 1.6, "HRMS"),
        (4.0, 1.6, "Projects"),
        (7.0, 1.6, "CRM"),
        (10.0, 1.6, "Payroll"),
        (1.0, 5.5, "Recruitment"),
        (4.0, 5.5, "Attendance"),
        (7.0, 5.5, "Inventory"),
        (10.0, 5.5, "Finance"),
        (0.4, 3.4, "Helpdesk"),
        (11.5, 3.4, "Reports"),
        (2.5, 2.5, "Analytics"),
        (9.4, 2.5, "Automation"),
    ]
    for lx, ty, label in nodes:
        c = card(s, Inches(lx), Inches(ty), Inches(1.9), Inches(0.7), LIGHT)
        add_text(c, label, size=13, bold=True, color=SECONDARY, align=PP_ALIGN.CENTER)
    note = s.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(12), Inches(0.4))
    add_text(note, "Shared data model · Unified permissions · Cross-module workflows", size=14, color=MUTED, align=PP_ALIGN.CENTER)

    def screenshot_slide(title: str, subtitle: str, shot_key: str, bullets: list[str]):
        s = blank_slide(prs)
        add_bg(s)
        title_block(s, title, subtitle)
        img_path = shots.get(shot_key) or shots.get("01-dashboard")
        s.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.45), Inches(8.6), Inches(5.5))
        panel = card(s, Inches(9.35), Inches(1.55), Inches(3.5), Inches(5.3), LIGHT)
        hb = s.shapes.add_textbox(Inches(9.55), Inches(1.8), Inches(3.1), Inches(0.4))
        add_text(hb, "Highlights", size=16, bold=True, color=PRIMARY)
        body = s.shapes.add_textbox(Inches(9.55), Inches(2.35), Inches(3.1), Inches(4.2))
        tf = add_text(body, "•  " + bullets[0], size=13, color=SLATE)
        for b in bullets[1:]:
            add_para(tf, "•  " + b, size=13, color=SLATE, space_before=10)
        return s

    # ---------- Slides 6-15 product ----------
    screenshot_slide(
        "Executive Dashboard",
        "Command center for the entire organization",
        "01-dashboard",
        ["KPI cards", "Revenue & pipeline", "Employees & attendance", "Projects & tasks", "Quick actions", "Notifications"],
    )
    screenshot_slide(
        "HRMS Dashboard",
        "Workforce health at a glance",
        "02-hrms-dashboard",
        ["Employee count", "Attendance summary", "Leave requests", "Department distribution", "Upcoming birthdays", "Expiring documents"],
    )
    screenshot_slide(
        "Employee Management",
        "Complete employee lifecycle records",
        "03b-employee-profile",
        ["Employee profile", "Documents", "Assets assigned", "Reporting structure", "Departments", "Timeline history"],
    )
    # Use profile as secondary visual note via another slide section - keep attendance
    screenshot_slide(
        "Attendance Management",
        "Shifts, daily records, and monthly visibility",
        "05-attendance-summary",
        ["Daily attendance", "Monthly summary", "Shift management", "Corrections workflow", "Present / late / leave", "Manager visibility"],
    )
    screenshot_slide(
        "Recruitment",
        "From requisition to offer in one pipeline",
        "08-job-openings",
        ["Job openings", "Candidates", "Interview rounds", "Offer management", "Hiring analytics", "Career portal ready"],
    )
    screenshot_slide(
        "Task Management",
        "Kanban, timeline, and collaborative execution",
        "12-tasks-board",
        ["Kanban board", "Priorities & owners", "Task details", "Comments", "Timeline view", "Attachments & watchers"],
    )
    screenshot_slide(
        "Project Management",
        "Delivery control with milestones and health",
        "14b-project-detail",
        ["Project portfolio", "Milestones", "Members", "Budget tracking", "Progress %", "Timeline & Gantt"],
    )
    screenshot_slide(
        "Resource Planning",
        "Capacity, allocation, and utilization",
        "15-resources-capacity",
        ["Team capacity", "Utilization bars", "Allocations", "Workload views", "Planner timeline", "Forecasting"],
    )
    screenshot_slide(
        "Analytics",
        "Executive reporting across CRM and operations",
        "17-reports",
        ["Revenue trends", "Pipeline value", "Win rate", "Lead conversion", "Outstanding invoices", "Finance reports"],
    )
    screenshot_slide(
        "CRM Module",
        "Optional sales suite — fully integrated",
        "19-pipeline",
        ["Lead list", "Sales pipeline", "Customers", "Quotations", "Invoices & payments", "Available as optional module"],
    )

    # ---------- Slide 16 Security ----------
    s = blank_slide(prs)
    add_bg(s)
    title_block(s, "Enterprise Security", "Designed for controlled access and auditability")
    sec_items = [
        ("security", "Role Based Access"),
        ("report", "Audit Logs"),
        ("shield", "Encryption Ready"),
        ("cloud", "Backups Ready"),
        ("module", "Granular Permissions"),
        ("clock", "Activity Logs"),
    ]
    for i, (icon, label) in enumerate(sec_items):
        col, row = i % 3, i // 3
        lx, ty = 0.7 + col * 4.15, 1.7 + row * 2.6
        card(s, Inches(lx), Inches(ty), Inches(3.9), Inches(2.3), LIGHT)
        s.shapes.add_picture(str(icons[icon]), Inches(lx + 1.45), Inches(ty + 0.4), Inches(0.9), Inches(0.9))
        tb = s.shapes.add_textbox(Inches(lx + 0.2), Inches(ty + 1.5), Inches(3.5), Inches(0.5))
        add_text(tb, label, size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    # optional audit screenshot inset
    if "22-audit-logs" in shots:
        s.shapes.add_picture(str(shots["22-audit-logs"]), Inches(9.2), Inches(5.5), Inches(3.7), Inches(1.7))

    # ---------- Slide 17 Industries ----------
    s = blank_slide(prs)
    add_bg(s)
    title_block(s, "Industry Solutions", "Configurable for sector-specific operating models")
    industries = [
        "Healthcare",
        "Construction",
        "Education",
        "IT / Software",
        "Retail",
        "Manufacturing",
        "Consulting",
        "Real Estate",
    ]
    for i, name in enumerate(industries):
        col, row = i % 4, i // 4
        lx, ty = 0.6 + col * 3.15, 1.8 + row * 2.5
        card(s, Inches(lx), Inches(ty), Inches(2.95), Inches(2.1), LIGHT)
        s.shapes.add_picture(str(icons["module"]), Inches(lx + 1.0), Inches(ty + 0.35), Inches(0.85), Inches(0.85))
        tb = s.shapes.add_textbox(Inches(lx + 0.15), Inches(ty + 1.35), Inches(2.65), Inches(0.5))
        add_text(tb, name, size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # ---------- Slide 18 Implementation ----------
    s = blank_slide(prs)
    add_bg(s)
    title_block(s, "Implementation Process", "A structured path from discovery to adoption")
    steps = ["Requirements", "Customization", "Development", "Testing", "Deployment", "Training", "Support"]
    for i, step in enumerate(steps):
        lx = 0.45 + i * 1.85
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(lx + 0.35), Inches(3.0), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY if i % 2 == 0 else ACCENT
        circle.line.fill.background()
        num = s.shapes.add_textbox(Inches(lx + 0.35), Inches(3.12), Inches(0.7), Inches(0.5))
        add_text(num, str(i + 1), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        label = s.shapes.add_textbox(Inches(lx), Inches(3.9), Inches(1.6), Inches(0.8))
        add_text(label, step, size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(lx + 1.2), Inches(3.3), Inches(0.85), Inches(0.08))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
            line.line.fill.background()
    note = s.shapes.add_textbox(Inches(0.6), Inches(5.3), Inches(12), Inches(1.2))
    add_text(
        note,
        "Dedicated onboarding, data migration support, role-based training, and ongoing success management "
        "ensure NovaCRM becomes operational quickly without disrupting day-to-day business.",
        size=15,
        color=SLATE,
    )

    # ---------- Slide 19 Roadmap ----------
    s = blank_slide(prs)
    add_bg(s)
    title_block(s, "Future Roadmap", "Continuously expanding the platform surface")
    left = card(s, Inches(0.6), Inches(1.7), Inches(5.8), Inches(5.0), LIGHT)
    right = card(s, Inches(6.9), Inches(1.7), Inches(5.8), Inches(5.0), LIGHT)
    lt = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(5.2), Inches(0.5))
    add_text(lt, "Current", size=22, bold=True, color=PRIMARY)
    rt = s.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(5.2), Inches(0.5))
    add_text(rt, "Upcoming", size=22, bold=True, color=ACCENT)
    for i, item in enumerate(["HRMS", "Projects & Tasks", "CRM Suite", "Resource Planning", "Recruitment", "Analytics"]):
        tb = s.shapes.add_textbox(Inches(1.0), Inches(2.7 + i * 0.55), Inches(5), Inches(0.45))
        add_text(tb, "●  " + item, size=16, color=SLATE)
    for i, item in enumerate(["Inventory", "Finance Deepening", "Procurement", "Helpdesk", "AI Assistants", "Advanced Analytics"]):
        tb = s.shapes.add_textbox(Inches(7.3), Inches(2.7 + i * 0.55), Inches(5), Inches(0.45))
        add_text(tb, "○  " + item, size=16, color=SLATE)

    # ---------- Slide 20 Comparison ----------
    s = blank_slide(prs)
    add_bg(s)
    title_block(s, "Why Choose NovaCRM", "Traditional stacks vs a unified platform")
    # table header
    card(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(5.2), LIGHT)
    card(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(5.2), RGBColor(0xEF, 0xF6, 0xFF))
    h1 = s.shapes.add_textbox(Inches(0.9), Inches(1.85), Inches(5.4), Inches(0.5))
    add_text(h1, "Traditional Software", size=20, bold=True, color=MUTED)
    h2 = s.shapes.add_textbox(Inches(7.1), Inches(1.85), Inches(5.4), Inches(0.5))
    add_text(h2, "NovaCRM", size=20, bold=True, color=PRIMARY)
    left_rows = ["Multiple Systems", "Multiple Logins", "Higher Cost", "Data Duplication", "Fragmented Reporting"]
    right_rows = ["Single Platform", "Unified Login", "Integrated Modules", "Scalable Architecture", "Modern UI"]
    for i, (a, b) in enumerate(zip(left_rows, right_rows)):
        la = s.shapes.add_textbox(Inches(1.0), Inches(2.7 + i * 0.7), Inches(5.2), Inches(0.5))
        add_text(la, "✕  " + a, size=16, color=ROSE)
        rb = s.shapes.add_textbox(Inches(7.2), Inches(2.7 + i * 0.7), Inches(5.2), Inches(0.5))
        add_text(rb, "✓  " + b, size=16, color=GREEN)

    # ---------- Slide 21 Mobile ----------
    s = blank_slide(prs)
    add_bg(s)
    title_block(s, "Mobile Experience", "Responsive web experience across phone, tablet, and laptop")
    s.shapes.add_picture(str(devices), Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))

    # ---------- Slide 22 Benefits ----------
    s = blank_slide(prs)
    add_bg(s)
    title_block(s, "Client Benefits", "Measurable impact for operators and executives")
    benefits = [
        ("cost", "Lower Costs"),
        ("prod", "Higher Productivity"),
        ("central", "Centralized Operations"),
        ("auto", "Automation"),
        ("report", "Better Reporting"),
        ("scale", "Scalability"),
    ]
    for i, (icon, label) in enumerate(benefits):
        col, row = i % 3, i // 3
        lx, ty = 0.7 + col * 4.15, 1.7 + row * 2.6
        card(s, Inches(lx), Inches(ty), Inches(3.9), Inches(2.3), LIGHT)
        s.shapes.add_picture(str(icons[icon]), Inches(lx + 1.45), Inches(ty + 0.35), Inches(0.95), Inches(0.95))
        tb = s.shapes.add_textbox(Inches(lx + 0.2), Inches(ty + 1.5), Inches(3.5), Inches(0.5))
        add_text(tb, label, size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # ---------- Slide 23 Thank You ----------
    s = blank_slide(prs)
    s.shapes.add_picture(str(cover_bg), 0, 0, W, H)
    t = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8), Inches(1))
    add_text(t, "Thank You", size=54, bold=True, color=WHITE)
    t2 = s.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(8), Inches(0.5))
    add_text(t2, "NovaCRM — One Platform. Complete Business Control.", size=18, color=RGBColor(0xBF, 0xDB, 0xFE))
    contact = s.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(7), Inches(2))
    tf = add_text(contact, "www.novacrm.com", size=16, color=WHITE)
    add_para(tf, "hello@novacrm.com", size=16, color=WHITE, space_before=8)
    add_para(tf, "+91 98765 43210", size=16, color=WHITE, space_before=8)
    add_para(tf, "Schedule a live enterprise demo", size=14, color=RGBColor(0x93, 0xC5, 0xFD), space_before=14)
    s.shapes.add_picture(str(qr), Inches(10.4), Inches(4.4), Inches(2.1), Inches(2.1))

    # Extra polish slides using more screenshots if we want 23+ - we already have 23

    prs.save(str(PPTX_PATH))
    print(f"Wrote {PPTX_PATH}")
    return PPTX_PATH


def export_pdf(pptx_path: Path) -> Path | None:
    """Try LibreOffice, then PowerPoint COM."""
    # LibreOffice
    for soffice in [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "soffice",
    ]:
        try:
            import subprocess

            cmd = [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(OUT_DIR), str(pptx_path)]
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
            if PDF_PATH.exists() or (OUT_DIR / "NovaCRM_Presentation.pdf").exists():
                print(f"PDF via LibreOffice: {PDF_PATH}")
                return PDF_PATH
            print("LibreOffice output:", r.stdout, r.stderr)
        except Exception as e:
            print("LibreOffice failed:", e)

    # PowerPoint COM
    try:
        import win32com.client  # type: ignore

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        deck = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
        # 32 = ppSaveAsPDF
        deck.SaveAs(str(PDF_PATH), 32)
        deck.Close()
        powerpoint.Quit()
        print(f"PDF via PowerPoint: {PDF_PATH}")
        return PDF_PATH
    except Exception as e:
        print("PowerPoint COM failed:", e)

    # Fallback: simple multi-page image PDF via Pillow (screenshot of rendered slides not available)
    # Create a PDF that embeds key visuals as pages for delivery continuity
    try:
        pages = []
        for name in [
            "01-dashboard",
            "02-hrms-dashboard",
            "12-tasks-board",
            "14-projects",
            "17-reports",
            "19-pipeline",
        ]:
            p = SHOTS / f"{name}.png"
            if p.exists():
                im = Image.open(p).convert("RGB").resize((1920, 1080), Image.Resampling.LANCZOS)
                pages.append(im)
        if pages:
            pages[0].save(PDF_PATH, "PDF", save_all=True, append_images=pages[1:], resolution=150.0)
            print(f"Fallback visual PDF: {PDF_PATH}")
            return PDF_PATH
    except Exception as e:
        print("Fallback PDF failed:", e)
    return None


if __name__ == "__main__":
    path = build()
    export_pdf(path)
